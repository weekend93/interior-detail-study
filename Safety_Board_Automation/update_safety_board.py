"""
무재해·무사고 기록판 PPTX 자동 갱신 + 이미지 저장

매일 1회 실행하면:
  1. template/ 안의 원본 PPTX를 읽어(원본은 수정하지 않음),
  2. '달성일수'(730일 목표 중 경과일)와 '현재' 날짜를 오늘 날짜 기준으로 다시 계산해 채우고,
  3. 슬라이드를 PNG 이미지로 변환해 바탕화면(Desktop)에 저장합니다.

'달성일수'는 시작일부터 오늘까지의 실제 경과일수로 매번 다시 계산하므로,
스케줄러 실행을 하루 건너뛰어도 다음 실행 때 자동으로 맞는 값이 채워집니다.
"""

import json
import shutil
import subprocess
import sys
import tempfile
from datetime import date, datetime
from pathlib import Path

from pptx import Presentation

SCRIPT_DIR = Path(__file__).resolve().parent


def load_config() -> dict:
    with open(SCRIPT_DIR / "config.json", encoding="utf-8") as f:
        return json.load(f)


def get_desktop_path() -> Path:
    """바탕화면 경로를 찾는다. Windows에서 OneDrive로 바탕화면이 리디렉션된 경우도 처리."""
    if sys.platform == "win32":
        try:
            import winreg

            key = winreg.OpenKey(
                winreg.HKEY_CURRENT_USER,
                r"Software\Microsoft\Windows\CurrentVersion\Explorer\Shell Folders",
            )
            value, _ = winreg.QueryValueEx(key, "Desktop")
            path = Path(winreg.ExpandEnvironmentStrings(value) if "%" in value else value)
            if path.exists():
                return path
        except Exception:
            pass
    return Path.home() / "Desktop"


def resolve_output_dir(cfg: dict) -> Path:
    if cfg["output_dir"] == "desktop":
        return get_desktop_path()
    return Path(cfg["output_dir"]).expanduser()


def set_run_text(slide, shape_id: int, run_index: int, new_text: str, expect_digits: bool = True):
    shape = next((s for s in slide.shapes if s.shape_id == shape_id), None)
    if shape is None:
        raise RuntimeError(
            f"shape_id={shape_id} 를 찾을 수 없습니다. PPTX 템플릿 구조가 변경된 것 같습니다."
        )
    run = shape.text_frame.paragraphs[0].runs[run_index]
    if expect_digits and not run.text.strip().isdigit():
        raise RuntimeError(
            f"shape_id={shape_id} run[{run_index}] 의 기존 값 '{run.text}' 이 "
            "숫자가 아닙니다. 템플릿이 예상과 달라 안전을 위해 중단합니다."
        )
    run.text = new_text


def build_updated_pptx(cfg: dict, today: date) -> Path:
    template_path = SCRIPT_DIR / cfg["template_pptx"]
    start_date = datetime.strptime(cfg["start_date"], "%Y-%m-%d").date()
    achievement_days = (today - start_date).days

    prs = Presentation(template_path)
    slide = prs.slides[0]

    set_run_text(slide, shape_id=16, run_index=1, new_text=f"{achievement_days} ")  # 달성일수
    set_run_text(slide, shape_id=24, run_index=0, new_text=f"{today.year} ", expect_digits=False)  # 현재-년
    set_run_text(slide, shape_id=25, run_index=0, new_text=f"{today.month:02d}")  # 현재-월
    set_run_text(slide, shape_id=26, run_index=0, new_text=f"{today.day:02d}")  # 현재-일

    work_pptx = Path(tempfile.mkdtemp(prefix="safety_board_")) / "current.pptx"
    prs.save(work_pptx)
    return work_pptx


def export_with_powerpoint(pptx_path: Path, png_path: Path, scale: int) -> bool:
    if sys.platform != "win32":
        return False
    try:
        import win32com.client
    except ImportError:
        return False

    try:
        prs = Presentation(pptx_path)
        width_px = int(prs.slide_width / 914400 * 96 * scale)
        height_px = int(prs.slide_height / 914400 * 96 * scale)

        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        deck = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        deck.Slides(1).Export(str(png_path), "PNG", width_px, height_px)
        deck.Close()
        powerpoint.Quit()
        return png_path.exists()
    except Exception as e:
        print(f"[PowerPoint 내보내기 실패, LibreOffice로 재시도] {e}")
        return False


def find_soffice() -> str | None:
    candidate = shutil.which("soffice")
    if candidate:
        return candidate
    for guess in (
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ):
        if Path(guess).exists():
            return guess
    return None


def export_with_libreoffice(pptx_path: Path, png_path: Path) -> bool:
    soffice = find_soffice()
    if not soffice:
        return False
    out_dir = png_path.parent
    subprocess.run(
        [soffice, "--headless", "--convert-to", "png", "--outdir", str(out_dir), str(pptx_path)],
        check=True,
        timeout=120,
    )
    generated = out_dir / (pptx_path.stem + ".png")
    if generated.exists() and generated != png_path:
        generated.replace(png_path)
    return png_path.exists()


def main():
    cfg = load_config()
    today = date.today()

    work_pptx = build_updated_pptx(cfg, today)

    output_dir = resolve_output_dir(cfg)
    output_dir.mkdir(parents=True, exist_ok=True)
    png_path = output_dir / today.strftime(cfg["filename_pattern"])

    ok = export_with_powerpoint(work_pptx, png_path, cfg.get("export_scale", 2))
    if not ok:
        ok = export_with_libreoffice(work_pptx, png_path)

    shutil.rmtree(work_pptx.parent, ignore_errors=True)

    if not ok:
        print(
            "이미지 저장에 실패했습니다. PowerPoint 또는 LibreOffice가 설치되어 있는지 확인해 주세요."
        )
        sys.exit(1)

    print(f"완료: {png_path}")


if __name__ == "__main__":
    main()
